"""
Gerador de PowerPoint Premium v3.0
Engenharia reversa baseada nas apresentações reais do cliente.

Estrutura de slides replicada do cliente:
  1. Capa (nome do condomínio + data da assembleia)
  2. Previsão Orçamentária (título + período)
  3-N. Tabelas de grupos de despesa (1 ou 2 grupos por slide)
  N+1. Composição do Rateio + Taxa Ideal
  N+2. Gráfico de pizza por grupo (% do total)
  N+3. Gráfico comparativo de barras (se tiver dados de 2 anos)
  Último. Slide de encerramento
"""

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
from io import BytesIO
import os
from datetime import datetime


# ─── PALETA FIEL AO CLIENTE (azul/cinza claro, header preto) ────────────────
C_BG        = RGBColor(0xFF, 0xFF, 0xFF)   # fundo branco
C_BG_LIGHT  = RGBColor(0xDA, 0xE8, 0xF5)   # azul claro (linhas alternadas da tabela)
C_HEADER    = RGBColor(0x00, 0x00, 0x00)   # cabeçalho preto
C_HEADER_TXT= RGBColor(0xFF, 0xFF, 0xFF)   # texto do cabeçalho branco
C_SUBTOTAL  = RGBColor(0x1F, 0x4E, 0x79)   # azul escuro para subtotais
C_SUBTOT_TXT= RGBColor(0xFF, 0xFF, 0xFF)
C_DARK_TEXT = RGBColor(0x00, 0x00, 0x00)   # texto principal preto
C_MUTED     = RGBColor(0x59, 0x59, 0x59)   # cinza médio
C_ACCENT    = RGBColor(0x1A, 0x56, 0xDB)   # azul accent
C_RED_ALERT = RGBColor(0xC0, 0x00, 0x00)
C_GREEN_OK  = RGBColor(0x37, 0x86, 0x10)
C_BORDER    = RGBColor(0xCC, 0xCC, 0xCC)

# Topo da página tem gradiente azul/cinza (skyline do cliente)
C_TOPO_DARK = RGBColor(0x9B, 0xB8, 0xD3)
C_TOPO_LIGHT= RGBColor(0xCC, 0xD9, 0xE8)


def br(r, g, b): return (r/255, g/255, b/255)


class PowerPointGenerator:
    """Gerador de apresentações no padrão do cliente Suport Condomínios."""

    W = Inches(13.33)
    H = Inches(7.5)

    def __init__(self, analyzer, output_path: str = "Relatorio_Condominio.pptx"):
        self.analyzer    = analyzer
        self.output_path = output_path
        self.prs         = Presentation()
        self.prs.slide_width  = self.W
        self.prs.slide_height = self.H
        self.r = analyzer.analysis_results

    # ─── PRIMITIVOS ──────────────────────────────────────────────────────────
    def _blank(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _bg(self, slide, color=C_BG):
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = color

    def _rect(self, slide, l, t, w, h, color, line_color=None):
        s = slide.shapes.add_shape(1, l, t, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = color
        if line_color:
            s.line.color.rgb = line_color
            s.line.width = Pt(0.5)
        else:
            s.line.fill.background()
        return s

    def _tb(self, slide, text, l, t, w, h, size=11, bold=False,
            color=C_DARK_TEXT, align=PP_ALIGN.LEFT, italic=False, wrap=True):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        return tb

    def _fmt_brl(self, val: float) -> str:
        """Formata em R$ brasileiro: R$ 1.234,56"""
        if val == 0:
            return "R$ 0,00"
        return f"R$ {val:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")

    def _fmt_pct(self, val: float) -> str:
        return f"{val:.1f}%".replace(".", ",")

    # ─── FAIXA TOPO (skyline estilo cliente) ─────────────────────────────────
    def _add_top_banner(self, slide):
        """Adiciona a faixa de skyline no topo, igual ao cliente."""
        self._rect(slide, 0, 0, self.W, Inches(0.72), C_TOPO_LIGHT)
        # Linha inferior da faixa
        self._rect(slide, 0, Inches(0.68), self.W, Pt(1), C_TOPO_DARK)

    # ─── SLIDE 1 — CAPA ──────────────────────────────────────────────────────
    def _slide_capa(self):
        slide = self._blank()
        self._bg(slide, C_BG)

        r = self.r
        nome_cond = r.get('nome_condominio', 'CONDOMÍNIO')
        ano       = r.get('ano_proximo', datetime.now().year)
        data_asm  = r.get('data_assembleia', '')

        # Fundo azul claro no topo (skyline simplificado)
        self._rect(slide, 0, 0, self.W, Inches(3.8), C_TOPO_LIGHT)
        self._rect(slide, 0, Inches(3.75), self.W, Pt(2), C_TOPO_DARK)

        # Nome principal
        self._tb(slide, nome_cond,
                 Inches(0.8), Inches(4.4), Inches(6.5), Inches(1.8),
                 size=44, bold=True, color=C_DARK_TEXT, align=PP_ALIGN.LEFT)

        # Linha separadora vertical
        self._rect(slide, Inches(7.5), Inches(4.5), Pt(1.5), Inches(1.2), C_MUTED)

        # Assembleia e data
        if data_asm:
            self._tb(slide, f"ASSEMBLEIA GERAL ORDINÁRIA\n{data_asm}",
                     Inches(7.8), Inches(4.6), Inches(5.0), Inches(1.0),
                     size=16, bold=False, color=C_MUTED, align=PP_ALIGN.CENTER)
        else:
            self._tb(slide, f"PREVISÃO ORÇAMENTÁRIA {ano}",
                     Inches(7.8), Inches(4.6), Inches(5.0), Inches(0.6),
                     size=18, bold=False, color=C_MUTED, align=PP_ALIGN.CENTER)

        print("  ✓ Slide 1: Capa")

    # ─── SLIDE 2 — TÍTULO PREVISÃO ORÇAMENTÁRIA ──────────────────────────────
    def _slide_titulo_previsao(self):
        slide = self._blank()
        self._bg(slide, C_BG)
        self._add_top_banner(slide)

        ano = self.r.get('ano_proximo', datetime.now().year)

        self._tb(slide, "PREVISÃO ORÇAMENTÁRIA",
                 Inches(1), Inches(2.5), Inches(11), Inches(1.5),
                 size=52, bold=True, color=C_DARK_TEXT, align=PP_ALIGN.CENTER)
        self._tb(slide, f"Período: {ano}",
                 Inches(1), Inches(4.3), Inches(11), Inches(1.0),
                 size=32, bold=False, color=C_MUTED, align=PP_ALIGN.CENTER)

        print("  ✓ Slide 2: Título Previsão")

    # ─── SLIDES DE TABELA POR GRUPO ──────────────────────────────────────────
    def _slide_grupos(self):
        """
        Cria slides de tabela com os grupos de despesa.
        Agrupa 1-2 grupos por slide (igual ao cliente).
        Formato: cabeçalho preto, linhas alternadas azul claro/branco,
                 total em azul escuro.
        """
        grupos = self.r.get('grupos', [])
        total_geral = self.r.get('total_despesas', 0)

        # Agrupa em pares (como o cliente faz)
        i = 0
        slide_num = 3
        while i < len(grupos):
            # Decide quantos grupos cabem: se o grupo tem muitos itens, fica sozinho
            g1 = grupos[i]
            g2 = grupos[i+1] if i+1 < len(grupos) else None

            n_itens_g1 = len(g1['itens'])
            n_itens_g2 = len(g2['itens']) if g2 else 0
            total_itens = n_itens_g1 + (n_itens_g2 if g2 else 0)

            # Se juntos têm mais de 12 itens, separa
            if g2 and total_itens > 12:
                self._slide_tabela_grupos([g1], total_geral, slide_num)
                slide_num += 1
                i += 1
            elif g2:
                self._slide_tabela_grupos([g1, g2], total_geral, slide_num)
                slide_num += 1
                i += 2
            else:
                self._slide_tabela_grupos([g1], total_geral, slide_num)
                slide_num += 1
                i += 1

    def _slide_tabela_grupos(self, grupos_lista, total_geral, slide_num):
        """Renderiza 1 ou 2 grupos de despesa em um slide no estilo tabela do cliente."""
        slide = self._blank()
        self._bg(slide, C_BG)
        self._add_top_banner(slide)

        num_grupos = len(grupos_lista)

        # Colunas da tabela
        COL_DESC = Inches(7.2)
        COL_TOT  = Inches(1.8)
        COL_APT  = Inches(1.8)
        COL_PCT  = Inches(2.0)
        ROW_H    = Inches(0.38)
        HEAD_H   = Inches(0.42)
        X0       = Inches(0.25)
        TOTAL_W  = COL_DESC + COL_TOT + COL_APT + COL_PCT

        y = Inches(0.8)

        for gi, grupo in enumerate(grupos_lista):
            # ── Cabeçalho do grupo ───────────────────────────────────────────
            self._rect(slide, X0, y, TOTAL_W, HEAD_H, C_HEADER)
            self._tb(slide,
                     f"{grupo['numero']}. {grupo['nome']}",
                     X0 + Inches(0.12), y + Pt(4), COL_DESC, HEAD_H,
                     size=12, bold=True, color=C_HEADER_TXT)
            self._tb(slide, "Total",
                     X0 + COL_DESC, y + Pt(4), COL_TOT, HEAD_H,
                     size=12, bold=True, color=C_HEADER_TXT, align=PP_ALIGN.CENTER)
            self._tb(slide, "Por Apto",
                     X0 + COL_DESC + COL_TOT, y + Pt(4), COL_APT, HEAD_H,
                     size=12, bold=True, color=C_HEADER_TXT, align=PP_ALIGN.CENTER)
            self._tb(slide, "Percentual",
                     X0 + COL_DESC + COL_TOT + COL_APT, y + Pt(4), COL_PCT, HEAD_H,
                     size=12, bold=True, color=C_HEADER_TXT, align=PP_ALIGN.CENTER)
            y += HEAD_H

            # ── Itens ────────────────────────────────────────────────────────
            # Percentual do grupo fica só na primeira linha
            pct_str = self._fmt_pct(grupo.get('percentual', 0)) + " do total"
            pct_shown = False

            for ii, item in enumerate(grupo['itens']):
                bg_cor = C_BG_LIGHT if ii % 2 == 0 else C_BG
                self._rect(slide, X0, y, TOTAL_W, ROW_H, bg_cor, C_BORDER)

                # Nome
                nome = item['nome']
                # Remove numeração redundante se já está no nome
                self._tb(slide, nome,
                         X0 + Inches(0.12), y + Pt(3), COL_DESC - Inches(0.15), ROW_H,
                         size=11, color=C_DARK_TEXT)
                # Total
                self._tb(slide, self._fmt_brl(item['total']),
                         X0 + COL_DESC, y + Pt(3), COL_TOT, ROW_H,
                         size=11, color=C_DARK_TEXT, align=PP_ALIGN.CENTER)
                # Por Apto
                por_apto = item.get('por_apto', 0)
                if por_apto == 0 and self.r.get('num_unidades', 0) > 0:
                    por_apto = item['total'] / self.r['num_unidades'] / 12
                self._tb(slide, self._fmt_brl(por_apto),
                         X0 + COL_DESC + COL_TOT, y + Pt(3), COL_APT, ROW_H,
                         size=11, color=C_DARK_TEXT, align=PP_ALIGN.CENTER)
                # Percentual (só na primeira linha do grupo)
                if not pct_shown:
                    self._tb(slide, pct_str,
                             X0 + COL_DESC + COL_TOT + COL_APT, y + Pt(3), COL_PCT, ROW_H,
                             size=10, color=C_MUTED, align=PP_ALIGN.CENTER, italic=True)
                    pct_shown = True

                y += ROW_H

            # ── Linha de TOTAL do grupo ──────────────────────────────────────
            self._rect(slide, X0, y, TOTAL_W, ROW_H, C_DARK_TEXT)
            por_apto_total = grupo['total'] / self.r.get('num_unidades', 202) / 12
            self._tb(slide, f"TOTAL - item {grupo['numero']}",
                     X0 + Inches(0.12), y + Pt(3), COL_DESC, ROW_H,
                     size=12, bold=True, color=C_HEADER_TXT)
            self._tb(slide, self._fmt_brl(grupo['total']),
                     X0 + COL_DESC, y + Pt(3), COL_TOT, ROW_H,
                     size=12, bold=True, color=C_HEADER_TXT, align=PP_ALIGN.CENTER)
            self._tb(slide, self._fmt_brl(por_apto_total),
                     X0 + COL_DESC + COL_TOT, y + Pt(3), COL_APT, ROW_H,
                     size=12, bold=True, color=C_HEADER_TXT, align=PP_ALIGN.CENTER)
            y += ROW_H + Inches(0.15)  # espaço entre grupos

        nomes = " + ".join(f"{g['numero']}. {g['nome']}" for g in grupos_lista)
        print(f"  ✓ Slide {slide_num}: Tabela — {nomes}")

    # ─── SLIDE — COMPOSIÇÃO DO RATEIO ────────────────────────────────────────
    def _slide_rateio(self):
        slide = self._blank()
        self._bg(slide, C_BG)
        self._add_top_banner(slide)

        r = self.r
        num_unid = r.get('num_unidades', 202)
        ano      = r.get('ano_proximo', datetime.now().year)

        X0 = Inches(0.8)
        W  = Inches(11.7)
        RH = Inches(0.55)
        y  = Inches(0.95)

        # ── TABELA 1: Composição do Rateio ───────────────────────────────────
        self._rect(slide, X0, y, W, RH, C_HEADER)
        self._tb(slide, "COMPOSIÇÃO DO RATEIO DAS DESPESAS",
                 X0 + Inches(0.1), y + Pt(6), W - Inches(0.2), RH,
                 size=14, bold=True, color=C_HEADER_TXT, align=PP_ALIGN.CENTER)
        y += RH

        linhas_rateio = [
            ("TOTAL DAS DESPESAS",          r['total_despesas'],          False),
            (f"FUNDO DE RESERVA (BC {r['fundo_reserva_pct']:.0f}%)", r['fundo_reserva'], False),
            (f"GARANTIDORA (BC {r['garantidora_pct']:.0f}%)",          r['garantidora'],  False),
            ("TOTAL DAS DESPESAS A SEREM RATEADAS", r['total_rateado'], True),
        ]
        for ii, (label, val, negrito) in enumerate(linhas_rateio):
            bg = C_BG_LIGHT if ii % 2 == 0 else C_BG
            self._rect(slide, X0, y, W, RH, bg, C_BORDER)
            self._tb(slide, label,
                     X0 + Inches(0.3), y + Pt(8), W * 0.6, RH,
                     size=13, bold=negrito, color=C_DARK_TEXT, align=PP_ALIGN.CENTER)
            self._tb(slide, self._fmt_brl(val),
                     X0 + W * 0.65, y + Pt(8), W * 0.32, RH,
                     size=13, bold=negrito, color=C_DARK_TEXT, align=PP_ALIGN.CENTER)
            y += RH

        y += Inches(0.35)

        # ── TABELA 2: Taxa Ideal ─────────────────────────────────────────────
        self._rect(slide, X0, y, W, RH, C_HEADER)
        self._tb(slide, f"COMPOSIÇÃO DA TAXA IDEAL - PREVISÃO ANUAL",
                 X0 + Inches(0.1), y + Pt(6), W - Inches(0.2), RH,
                 size=14, bold=True, color=C_HEADER_TXT, align=PP_ALIGN.CENTER)
        y += RH

        taxa_mensal = r.get('taxa_ideal_mensal', 0)
        taxa_fundo  = r.get('taxa_fundo_reserva', 0)
        taxa_total  = taxa_mensal  # já inclui o fundo? Verifica
        # No cliente: "Despesas Mensais" = taxa sem fundo, depois soma fundo
        taxa_desp_pura = r['total_despesas'] / num_unid / 12
        taxa_fundo_un  = r['fundo_reserva'] / num_unid / 12
        taxa_total_real = taxa_desp_pura + taxa_fundo_un

        linhas_taxa = [
            (f"DESPESAS MENSAIS - {num_unid} UNIDADES", taxa_desp_pura, False),
            (f"FUNDO DE RESERVA (BC {r['fundo_reserva_pct']:.0f}%)", taxa_fundo_un, False),
            ("TOTAL DA TAXA DE CONDOMÍNIO",              taxa_total_real, True),
        ]
        for ii, (label, val, negrito) in enumerate(linhas_taxa):
            bg = C_BG_LIGHT if ii % 2 == 0 else C_BG
            self._rect(slide, X0, y, W, RH, bg, C_BORDER)
            self._tb(slide, label,
                     X0 + Inches(0.3), y + Pt(8), W * 0.6, RH,
                     size=13, bold=negrito, color=C_DARK_TEXT, align=PP_ALIGN.CENTER)
            self._tb(slide, self._fmt_brl(val),
                     X0 + W * 0.65, y + Pt(8), W * 0.32, RH,
                     size=13, bold=negrito, color=C_DARK_TEXT, align=PP_ALIGN.CENTER)
            y += RH

        y += Inches(0.45)

        # ── Linha de reajuste ─────────────────────────────────────────────────
        reajuste = r.get('reajuste_pct')
        taxa_atual = r.get('taxa_atual')
        if reajuste is not None and taxa_atual:
            msg = (f"REAJUSTE DE {reajuste:.2f}% EM RELAÇÃO AO VALOR ATUAL, "
                   f"DE {self._fmt_brl(taxa_atual)}.")
            self._rect(slide, X0, y, W, Inches(0.65), C_BG_LIGHT, C_BORDER)
            self._tb(slide, msg,
                     X0 + Inches(0.2), y + Pt(10), W - Inches(0.4), Inches(0.65),
                     size=13, bold=False, color=C_DARK_TEXT, align=PP_ALIGN.LEFT)

        print("  ✓ Slide Rateio: Composição e Taxa Ideal")

    # ─── SLIDE — GRÁFICO PIZZA POR GRUPO ─────────────────────────────────────
    def _slide_grafico_pizza(self):
        slide = self._blank()
        self._bg(slide, C_BG)
        self._add_top_banner(slide)

        grupos = self.r.get('grupos', [])
        if not grupos:
            return

        self._tb(slide, "DISTRIBUIÇÃO DAS DESPESAS POR CATEGORIA",
                 Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.55),
                 size=20, bold=True, color=C_DARK_TEXT, align=PP_ALIGN.CENTER)

        nomes  = [f"{g['numero']}. {g['nome']}" for g in grupos]
        totais = [g['total'] for g in grupos]
        total_geral = sum(totais)

        # Cores harmônicas com o estilo do cliente
        cores = [
            br(0x1A, 0x56, 0xDB),  # azul
            br(0x0E, 0x9F, 0x6E),  # verde
            br(0xF5, 0xA6, 0x23),  # âmbar
            br(0xE0, 0x2A, 0x2A),  # vermelho
            br(0x6D, 0x28, 0xD9),  # roxo
            br(0x06, 0x82, 0x90),  # teal
            br(0x9B, 0x59, 0xB6),  # lilás
            br(0x27, 0xAE, 0x60),  # verde escuro
        ]
        cores = (cores * 4)[:len(grupos)]

        fig, ax = plt.subplots(figsize=(7, 5.5))
        fig.patch.set_facecolor('white')

        wedges, texts, autotexts = ax.pie(
            totais,
            labels=None,
            autopct=lambda p: f'{p:.1f}%' if p > 3 else '',
            colors=cores,
            startangle=140,
            pctdistance=0.75,
            wedgeprops={'linewidth': 1.5, 'edgecolor': 'white'},
        )
        for at in autotexts:
            at.set_fontsize(9)
            at.set_fontweight('bold')
            at.set_color('white')

        # Legenda
        labels_leg = [f"{g['numero']}. {g['nome'][:28]} — {self._fmt_brl(g['total'])} ({g['total']/total_geral*100:.1f}%)"
                      for g in grupos]
        ax.legend(wedges, labels_leg,
                  loc='center left', bbox_to_anchor=(1.02, 0.5),
                  fontsize=8.5, frameon=False)

        ax.set_aspect('equal')
        plt.tight_layout(pad=0.5)

        buf = BytesIO()
        fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                    facecolor='white', transparent=False)
        buf.seek(0)
        plt.close(fig)

        slide.shapes.add_picture(buf, Inches(0.3), Inches(1.5),
                                  width=Inches(12.5))
        print("  ✓ Slide Gráfico Pizza: Distribuição por Categoria")

    # ─── SLIDE — GRÁFICO BARRAS POR GRUPO ────────────────────────────────────
    def _slide_grafico_barras(self):
        slide = self._blank()
        self._bg(slide, C_BG)
        self._add_top_banner(slide)

        grupos = self.r.get('grupos', [])
        if not grupos:
            return

        self._tb(slide, "VALOR POR CATEGORIA DE DESPESA",
                 Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.55),
                 size=20, bold=True, color=C_DARK_TEXT, align=PP_ALIGN.CENTER)

        nomes  = [f"{g['numero']}. {g['nome']}" for g in grupos]
        totais = [g['total'] / 1000 for g in grupos]  # em mil
        pcts   = [g.get('percentual', 0) for g in grupos]

        cores = [
            br(0x1A, 0x56, 0xDB),
            br(0x0E, 0x9F, 0x6E),
            br(0xF5, 0xA6, 0x23),
            br(0xE0, 0x2A, 0x2A),
            br(0x6D, 0x28, 0xD9),
            br(0x06, 0x82, 0x90),
            br(0x9B, 0x59, 0xB6),
            br(0x27, 0xAE, 0x60),
        ]
        cores = (cores * 4)[:len(grupos)]

        short_names = [n[:22] + '…' if len(n) > 22 else n for n in nomes]

        fig, ax = plt.subplots(figsize=(12.5, 4.8))
        fig.patch.set_facecolor('white')
        ax.set_facecolor('#F8FAFF')

        x = np.arange(len(grupos))
        bars = ax.bar(x, totais, color=cores, width=0.55,
                      linewidth=0, alpha=0.92, zorder=3)

        # Anotações: valor + %
        for bar, val, pct in zip(bars, totais, pcts):
            ax.text(bar.get_x() + bar.get_width()/2,
                    bar.get_height() + 0.8,
                    f'R${val:.0f}K\n({pct:.1f}%)',
                    ha='center', va='bottom', fontsize=8.5,
                    fontweight='bold', color='#1E293B')

        ax.set_xticks(x)
        ax.set_xticklabels(short_names, rotation=30, ha='right', fontsize=9.5)
        ax.set_ylabel('Valor (R$ mil)', fontsize=11)
        ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v, _: f'R${v:.0f}K'))
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.spines['left'].set_visible(False)
        ax.grid(axis='y', color='#E2E8F0', linewidth=0.8)
        ax.tick_params(left=False)
        ax.set_axisbelow(True)

        plt.tight_layout(pad=1.0)

        buf = BytesIO()
        fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                    facecolor='white', transparent=False)
        buf.seek(0)
        plt.close(fig)

        slide.shapes.add_picture(buf, Inches(0.3), Inches(1.45),
                                  width=Inches(12.7))
        print("  ✓ Slide Gráfico Barras: Valor por Categoria")

    # ─── SLIDE — ENCERRAMENTO ────────────────────────────────────────────────
    def _slide_encerramento(self):
        slide = self._blank()
        self._bg(slide, C_BG)
        self._add_top_banner(slide)

        r = self.r
        reajuste = r.get('reajuste_pct')
        taxa_atual = r.get('taxa_atual')
        taxa_nova  = r.get('taxa_ideal_mensal', 0) + r.get('taxa_fundo_reserva', 0)

        # Cálculo correto da taxa total
        num_unid = r.get('num_unidades', 202)
        taxa_desp  = r['total_despesas'] / num_unid / 12
        taxa_fundo = r['fundo_reserva'] / num_unid / 12
        taxa_nova  = taxa_desp + taxa_fundo

        if reajuste is not None and reajuste > 0:
            cor_status = C_RED_ALERT
            status_txt = f"Reajuste necessário de {reajuste:.2f}%"
            icon = "⚠"
        elif reajuste is not None:
            cor_status = C_GREEN_OK
            status_txt = "Taxa atual é suficiente para cobrir as despesas"
            icon = "✔"
        else:
            cor_status = C_DARK_TEXT
            status_txt = "Análise financeira concluída"
            icon = "✔"

        self._tb(slide, "RESUMO FINANCEIRO",
                 Inches(1), Inches(1.0), Inches(11.3), Inches(0.7),
                 size=30, bold=True, color=C_DARK_TEXT, align=PP_ALIGN.CENTER)

        # Cards de resumo
        cards = [
            ("Total das Despesas",    self._fmt_brl(r['total_despesas'])),
            ("Total a Ratear",        self._fmt_brl(r['total_rateado'])),
            ("Taxa Ideal / mês",      self._fmt_brl(taxa_nova)),
        ]
        if taxa_atual:
            cards.append(("Taxa Atual / mês", self._fmt_brl(taxa_atual)))

        card_w = Inches(2.8)
        card_h = Inches(1.4)
        spacing = (self.W - card_w * len(cards)) / (len(cards) + 1)
        y_card  = Inches(1.9)

        for ci, (label, val) in enumerate(cards):
            cx = spacing + ci * (card_w + spacing)
            self._rect(slide, cx, y_card, card_w, card_h,
                       C_BG_LIGHT, C_BORDER)
            self._tb(slide, label,
                     cx + Inches(0.1), y_card + Inches(0.1),
                     card_w - Inches(0.2), Inches(0.45),
                     size=10, color=C_MUTED, align=PP_ALIGN.CENTER)
            self._tb(slide, val,
                     cx + Inches(0.05), y_card + Inches(0.55),
                     card_w - Inches(0.1), Inches(0.7),
                     size=18, bold=True, color=C_DARK_TEXT, align=PP_ALIGN.CENTER)

        # Status
        self._rect(slide, Inches(1.5), Inches(3.7), Inches(10.3), Inches(0.85),
                   cor_status)
        self._tb(slide, f"{icon}  {status_txt}",
                 Inches(1.7), Inches(3.82), Inches(10.0), Inches(0.65),
                 size=16, bold=True, color=C_HEADER_TXT, align=PP_ALIGN.CENTER)

        # Linha informativa de reajuste
        if reajuste is not None and taxa_atual:
            info = (f"A taxa sugerida de {self._fmt_brl(taxa_nova)}/mês representa "
                    f"um reajuste de {reajuste:.2f}% em relação à taxa atual de "
                    f"{self._fmt_brl(taxa_atual)}.")
            self._tb(slide, info,
                     Inches(1.5), Inches(4.75), Inches(10.3), Inches(0.7),
                     size=12, color=C_MUTED, align=PP_ALIGN.CENTER)

        # Rodapé
        self._tb(slide,
                 f"Relatório gerado em {datetime.now().strftime('%d/%m/%Y')} · Gestor Financeiro de Condomínios v3.0",
                 Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4),
                 size=9, color=C_MUTED, align=PP_ALIGN.CENTER)

        print("  ✓ Slide Encerramento: Resumo Financeiro")

    # ─── GENERATE ────────────────────────────────────────────────────────────
    def generate(self):
        print("\nGerando apresentação PowerPoint (estilo cliente)...")

        self._slide_capa()
        self._slide_titulo_previsao()
        self._slide_grupos()
        self._slide_rateio()
        self._slide_grafico_pizza()
        self._slide_grafico_barras()
        self._slide_encerramento()

        self.prs.save(self.output_path)
        print(f"\n✓ Apresentação salva: {self.output_path}")
        return self.output_path
